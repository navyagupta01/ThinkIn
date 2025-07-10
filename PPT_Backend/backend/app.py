from flask_cors import CORS
from flask import Flask, request, send_file, jsonify
import requests, os, io, zipfile, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import seaborn as sns
import pandas as pd
import numpy as np
from datetime import datetime
import json

app = Flask(__name__)
CORS(app)

# Configuration
API_KEY = os.getenv("OPENROUTER_API_KEY", "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")
MODEL = "gpt-4o-mini"

def call_openrouter(prompt: str) -> str:
    """Make API call to OpenRouter"""
    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "stream": False
    }
    
    print("Request payload:", data)
    resp = requests.post(url, headers=headers, json=data)
    
    if resp.status_code != 200:
        raise ValueError(f"OpenRouter API error: {resp.status_code} - {resp.text}")
    
    try:
        resp_json = resp.json()
        print("API response:", resp_json)
        return resp_json["choices"][0]["message"]["content"]
    except Exception as e:
        print("Error parsing response:", e)
        raise ValueError(f"Unexpected API response structure: {resp.text}")

def sanitize_filename(filename):
    """Sanitize filename to remove invalid characters"""
    filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
    filename = re.sub(r'[^\w\s\-_.]', '', filename)
    filename = filename.strip()
    if len(filename) > 100:
        filename = filename[:100]
    return filename

def gen_structured_outline(topic, sessions, details, syllabus):
    """Generate structured outline for multiple sessions with detailed content"""
    if syllabus.strip():
        prompt = f"""
        Create a comprehensive curriculum for teaching '{topic}' across {sessions} sessions, strictly based on the provided syllabus.

        Syllabus: {syllabus}

        IMPORTANT REQUIREMENTS:
        1. Each bullet point should be 15-25 words long with detailed explanations
        2. Include 8-12 slides per session minimum
        3. Make content educational and comprehensive
        4. For code slides, provide complete, working examples with explanations
        5. Include practical examples and real-world applications
        6. Distribute the syllabus topics evenly across the specified number of sessions
        7. Each session should cover specific topics from the syllabus

        Format your response as JSON with the following structure:
        {{
            "course_title": "Main Course Title",
            "sessions": [
                {{
                    "session_number": 1,
                    "session_title": "Session Title (covering specific syllabus topics)",
                    "slides": [
                        {{
                            "slide_title": "Slide Title",
                            "content": [
                                "Detailed bullet point explaining key concept from syllabus with examples (15-25 words)",
                                "Another comprehensive point covering syllabus topic (15-25 words)",
                                "Third point with practical applications from syllabus (15-25 words)",
                                "Fourth point covering advanced syllabus concepts (15-25 words)"
                            ],
                            "slide_type": "content"
                        }},
                        {{
                            "slide_title": "Code Example: [Specific Syllabus Topic]",
                            "content": [
                                "This code demonstrates how to implement [specific syllabus functionality]",
                                "Key points to understand in this implementation",
                                "Common use cases and practical applications"
                            ],
                            "slide_type": "code",
                            "code_content": [
                                "# Complete working code example",
                                "# with detailed comments explaining each step",
                                "def example_function():",
                                "    # Implementation details",
                                "    return result"
                            ]
                        }},
                        {{
                            "slide_title": "Chart Title: [Syllabus Data Description]",
                            "content": [
                                "This chart shows relationships in {topic} from syllabus",
                                "Analysis of trends and patterns from syllabus topics",
                                "Implications for practical applications"
                            ],
                            "slide_type": "chart",
                            "chart_data": {{
                                "chart_type": "bar",
                                "data": [
                                    ["Category", "Value"],
                                    ["Item A", 45],
                                    ["Item B", 67],
                                    ["Item C", 32],
                                    ["Item D", 78]
                                ]
                            }}
                        }}
                    ]
                }}
            ]
        }}

        Make sure to:
        - Create substantial, educational content for each slide based on syllabus topics
        - Include proper code examples with comments when relevant to syllabus
        - Provide meaningful chart data related to syllabus topics
        - Ensure each bullet point is informative and detailed
        - Cover both theoretical concepts and practical applications from syllabus
        - Align session content with specific syllabus topics

        Topic: {topic}
        Additional details: {details}
        """
    else:
        prompt = f"""
        Create a comprehensive curriculum for teaching '{topic}' across {sessions} sessions.

        IMPORTANT REQUIREMENTS:
        1. Each bullet point should be 15-25 words long with detailed explanations
        2. Include 8-12 slides per session minimum
        3. Make content educational and comprehensive
        4. For code slides, provide complete, working examples with explanations
        5. Include practical examples and real-world applications

        Format your response as JSON with the following structure:
        {{
            "course_title": "Main Course Title",
            "sessions": [
                {{
                    "session_number": 1,
                    "session_title": "Session Title",
                    "slides": [
                        {{
                            "slide_title": "Slide Title",
                            "content": [
                                "Detailed bullet point explaining key concept with examples and context (15-25 words)",
                                "Another comprehensive point covering important aspects of the topic (15-25 words)",
                                "Third point with practical applications and real-world usage (15-25 words)",
                                "Fourth point covering advanced concepts or common pitfalls (15-25 words)"
                            ],
                            "slide_type": "content"
                        }},
                        {{
                            "slide_title": "Code Example: [Specific Topic]",
                            "content": [
                                "This code demonstrates how to implement [specific functionality] in {topic}",
                                "Key points to understand in this implementation",
                                "Common use cases and practical applications"
                            ],
                            "slide_type": "code",
                            "code_content": [
                                "# Complete working code example",
                                "# with detailed comments explaining each step",
                                "def example_function():",
                                "    # Implementation details",
                                "    return result"
                            ]
                        }},
                        {{
                            "slide_title": "Chart Title: [Data Description]",
                            "content": [
                                "This chart shows the relationship between key metrics in {topic}",
                                "Analysis of trends and patterns visible in the data",
                                "Implications for practical applications"
                            ],
                            "slide_type": "chart",
                            "chart_data": {{
                                "chart_type": "bar",
                                "data": [
                                    ["Category", "Value"],
                                    ["Item A", 45],
                                    ["Item B", 67],
                                    ["Item C", 32],
                                    ["Item D", 78]
                                ]
                            }}
                        }}
                    ]
                }}
            ]
        }}

        Make sure to:
        - Create substantial, educational content for each slide
        - Include proper code examples with comments when relevant
        - Provide meaningful chart data related to the topic
        - Ensure each bullet point is informative and detailed
        - Cover both theoretical concepts and practical applications

        Topic: {topic}
        Additional details: {details}
        """
    return call_openrouter(prompt)

def parse_text_outline(outline_text, topic, sessions, syllabus):
    """Parse plain text outline into structured format as fallback with enhanced content"""
    lines = outline_text.split('\n')
    sessions_data = []
    current_session = None
    current_slide = None

    # If syllabus is provided, split it into topics
    syllabus_topics = [t.strip() for t in syllabus.split(',')] if syllabus.strip() else []
    if syllabus_topics:
        sessions_per_topic = max(1, sessions // len(syllabus_topics))
        remaining_sessions = sessions % len(syllabus_topics)
    else:
        sessions_per_topic = sessions
        remaining_sessions = 0

    session_counter = 0
    topic_index = 0

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Session detection
        if line.startswith('### Session') or line.startswith('Session'):
            if current_session:
                sessions_data.append(current_session)

            if syllabus_topics and topic_index < len(syllabus_topics):
                session_title = f"Session {session_counter + 1}: {syllabus_topics[topic_index]}"
                if session_counter % sessions_per_topic == 0 and topic_index < len(syllabus_topics) - 1:
                    topic_index += 1
            else:
                session_title = line.replace('### ', '').replace('**', '')

            current_session = {
                "session_number": len(sessions_data) + 1,
                "session_title": session_title,
                "slides": []
            }
            current_slide = None
            session_counter += 1

        # Slide detection
        elif line.startswith('**Slide') and current_session:
            if current_slide:
                current_session['slides'].append(current_slide)

            slide_title = line.replace('**', '').replace('Slide ', '').split(':')[1].strip() if ':' in line else line.replace('**', '')
            if syllabus_topics and topic_index < len(syllabus_topics):
                slide_title = f"{slide_title} ({syllabus_topics[topic_index]})"

            slide_type = "content"
            chart_data = None
            code_content = None

            if any(keyword in slide_title.lower() for keyword in ["code", "example", "implementation", "syntax", "programming"]):
                slide_type = "code"
                code_content = [
                    f"# {topic} - Code Example",
                    "# This example demonstrates key concepts",
                    "def main_function():",
                    "    # Implementation details",
                    "    result = process_data()",
                    "    return result",
                    "",
                    "# Usage example",
                    "output = main_function()",
                    "print(output)"
                ]
            elif any(keyword in slide_title.lower() for keyword in ["statistics", "data", "comparison", "trend", "growth", "analysis", "metrics", "performance", "survey", "timeline"]):
                slide_type = "chart"
                chart_data = {
                    "chart_type": "bar",
                    "data": [
                        ["Category", "Value"],
                        ["Beginner Level", 35],
                        ["Intermediate Level", 45],
                        ["Advanced Level", 25],
                        ["Expert Level", 15]
                    ]
                }

            current_slide = {
                "slide_title": slide_title,
                "content": [
                    f"Comprehensive overview of {slide_title.lower()} and its importance in {topic}",
                    f"Key principles and best practices for implementing {slide_title.lower()} effectively",
                    f"Common challenges and solutions when working with {slide_title.lower()}",
                    f"Real-world applications and practical examples of {slide_title.lower()}"
                ],
                "slide_type": slide_type
            }

            if chart_data:
                current_slide["chart_data"] = chart_data
            if code_content:
                current_slide["code_content"] = code_content

        # Bullet points (enhance short ones)
        elif line.startswith('-') and current_slide:
            bullet = line[1:].strip()
            if len(bullet.split()) < 8:
                bullet = f"{bullet} - detailed explanation with practical examples and implementation considerations"
            current_slide['content'].append(bullet)

        # Code content
        elif current_slide and any(keyword in line for keyword in ['import ', 'def ', 'class ', 'model =', 'print(']):
            if 'code_content' not in current_slide:
                current_slide['code_content'] = []
            current_slide['code_content'].append(line)

    # Add the last session and slide
    if current_slide and current_session:
        current_session['slides'].append(current_slide)
    if current_session:
        sessions_data.append(current_session)

    # Ensure enough sessions for syllabus topics
    while len(sessions_data) < sessions:
        session_num = len(sessions_data) + 1
        topic_title = syllabus_topics[topic_index] if topic_index < len(syllabus_topics) else f"Additional Topics for {topic}"
        if topic_index < len(syllabus_topics) and session_num % sessions_per_topic == 0:
            topic_index += 1
        sessions_data.append({
            "session_number": session_num,
            "session_title": f"Session {session_num}: {topic_title}",
            "slides": [
                {
                    "slide_title": f"Introduction to {topic_title}",
                    "content": [
                        f"Overview of {topic_title} and its significance in {topic}",
                        f"Key concepts and principles related to {topic_title}",
                        f"Practical applications of {topic_title} in real-world scenarios",
                        f"Best practices for understanding and applying {topic_title}"
                    ],
                    "slide_type": "content"
                }
            ]
        })

    # Ensure we have enough content per session
    for session in sessions_data:
        while len(session['slides']) < 8:
            slide_num = len(session['slides']) + 1
            session['slides'].append({
                "slide_title": f"Additional Topic {slide_num}",
                "content": [
                    f"Important concepts and principles related to {topic} that enhance understanding",
                    f"Practical applications and real-world scenarios where these concepts are applied",
                    f"Best practices and common approaches used by professionals in the field",
                    f"Key takeaways and actionable insights for implementing these concepts effectively"
                ],
                "slide_type": "content"
            })

    return {
        "course_title": topic,
        "sessions": sessions_data
    }

def create_chart(chart_data, filename):
    """Create various types of charts with professional styling"""
    try:
        plt.figure(figsize=(12, 8))
        plt.style.use('default')
        
        chart_type = chart_data.get('chart_type', 'bar')
        data = chart_data.get('data', [])
        
        if len(data) < 2:
            return None
        
        headers = data[0]
        values = data[1:]
        
        colors = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6A994E', '#7209B7']
        
        if chart_type == 'bar':
            categories = [str(row[0]) for row in values]
            nums = []
            for row in values:
                try:
                    nums.append(float(row[1]))
                except (ValueError, IndexError):
                    nums.append(0)
            
            bars = plt.bar(categories, nums, color=colors[:len(categories)], 
                          edgecolor='white', linewidth=1.5, alpha=0.8)
            
            for bar, num in zip(bars, nums):
                plt.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(nums)*0.01,
                        f'{num:.1f}', ha='center', va='bottom', fontweight='bold')
            
            plt.title(f'{headers[0]} vs {headers[1]}', fontsize=18, fontweight='bold', pad=20)
            plt.xlabel(headers[0], fontsize=14, fontweight='bold')
            plt.ylabel(headers[1], fontsize=14, fontweight='bold')
            plt.xticks(rotation=45, ha='right')
            plt.grid(axis='y', alpha=0.3, linestyle='--')
            
        elif chart_type == 'line':
            x_vals = [str(row[0]) for row in values]
            y_vals = []
            for row in values:
                try:
                    y_vals.append(float(row[1]))
                except (ValueError, IndexError):
                    y_vals.append(0)
            
            plt.plot(x_vals, y_vals, marker='o', linewidth=3, markersize=10, 
                    color=colors[0], markerfacecolor=colors[1], markeredgecolor='white', 
                    markeredgewidth=2)
            
            for x, y in zip(x_vals, y_vals):
                plt.text(x, y + max(y_vals)*0.02, f'{y:.1f}', ha='center', va='bottom', 
                        fontweight='bold')
            
            plt.title(f'{headers[0]} vs {headers[1]}', fontsize=18, fontweight='bold', pad=20)
            plt.xlabel(headers[0], fontsize=14, fontweight='bold')
            plt.ylabel(headers[1], fontsize=14, fontweight='bold')
            plt.xticks(rotation=45, ha='right')
            plt.grid(True, alpha=0.3, linestyle='--')
            
        elif chart_type == 'pie':
            labels = [str(row[0]) for row in values]
            sizes = []
            for row in values:
                try:
                    sizes.append(float(row[1]))
                except (ValueError, IndexError):
                    sizes.append(1)
            
            wedges, texts, autotexts = plt.pie(sizes, labels=labels, autopct='%1.1f%%', 
                                             startangle=90, colors=colors[:len(labels)],
                                             explode=[0.05]*len(labels), shadow=True)
            
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontweight('bold')
                autotext.set_fontsize(12)
            
            for text in texts:
                text.set_fontsize(12)
                text.set_fontweight('bold')
            
            plt.title(f'Distribution of {headers[0]}', fontsize=18, fontweight='bold', pad=20)
        
        plt.gca().set_facecolor('#f8f9fa')
        plt.gcf().patch.set_facecolor('white')
        
        plt.tight_layout()
        plt.savefig(filename, dpi=300, bbox_inches='tight', format='png', 
                   facecolor='white', edgecolor='none')
        plt.close()
        return filename
        
    except Exception as e:
        print(f"Error creating chart: {e}")
        plt.close()
        return None

def create_ppt_session(session_data, topic):
    """Create a PowerPoint presentation for a single session with improved formatting"""
    try:
        prs = Presentation()
        
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = session_data.get('session_title', 'Session Title')
        
        if len(title_slide.shapes.placeholders) > 1:
            title_slide.shapes.placeholders[1].text = f"Course: {topic}"
        
        title_shape = title_slide.shapes.title
        if title_shape.has_text_frame:
            title_frame = title_shape.text_frame
            if title_frame.paragraphs:
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(36)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 51, 102)
        
        for slide_data in session_data.get('slides', []):
            slide_type = slide_data.get('slide_type', 'content')
            slide_title = slide_data.get('slide_title', 'Slide Title')
            
            if slide_type == 'chart' and 'chart_data' in slide_data:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = slide_title
                title_para.font.size = Pt(28)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 51, 102)
                
                content_items = slide_data.get('content', [])
                if content_items:
                    desc_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1))
                    desc_frame = desc_shape.text_frame
                    desc_para = desc_frame.paragraphs[0]
                    desc_para.text = content_items[0] if content_items else "Chart analysis and insights"
                    desc_para.font.size = Pt(16)
                    desc_para.font.color.rgb = RGBColor(51, 51, 51)
                
                chart_filename = f"temp_chart_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}.png"
                chart_path = create_chart(slide_data['chart_data'], chart_filename)
                
                if chart_path and os.path.exists(chart_path):
                    try:
                        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(2.5), 
                                               Inches(9), Inches(4.5))
                    except Exception as e:
                        print(f"Error adding chart to slide: {e}")
                    finally:
                        try:
                            os.remove(chart_path)
                        except:
                            pass
                else:
                    _create_content_slide(prs, slide_data)
                    
            elif slide_type == 'code':
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = slide_title
                title_para.font.size = Pt(28)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 51, 102)
                
                content_items = slide_data.get('content', [])
                if content_items:
                    exp_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
                    exp_frame = exp_shape.text_frame
                    exp_frame.word_wrap = True
                    
                    for i, item in enumerate(content_items[:2]):
                        if i == 0:
                            p = exp_frame.paragraphs[0]
                        else:
                            p = exp_frame.add_paragraph()
                        p.text = f"• {item}"
                        p.font.size = Pt(16)
                        p.font.color.rgb = RGBColor(51, 51, 51)
                        p.space_after = Pt(8)
                
                code_content = slide_data.get('code_content', [])
                if code_content:
                    code_shape = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(4))
                    code_frame = code_shape.text_frame
                    code_frame.word_wrap = True
                    code_frame.margin_left = Inches(0.2)
                    code_frame.margin_right = Inches(0.2)
                    code_frame.margin_top = Inches(0.2)
                    code_frame.margin_bottom = Inches(0.2)
                    
                    fill = code_shape.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(45, 45, 45)
                    
                    line = code_shape.line
                    line.color.rgb = RGBColor(100, 100, 100)
                    line.width = Pt(1)
                    
                    code_text = '\n'.join(str(line) for line in code_content)
                    code_para = code_frame.paragraphs[0]
                    code_para.text = code_text
                    code_para.font.name = 'Consolas'
                    code_para.font.size = Pt(14)
                    code_para.font.color.rgb = RGBColor(220, 220, 220)
                    code_para.space_after = Pt(0)
                    
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_title
                
                title_shape = slide.shapes.title
                if title_shape.has_text_frame:
                    title_frame = title_shape.text_frame
                    if title_frame.paragraphs:
                        title_para = title_frame.paragraphs[0]
                        title_para.font.size = Pt(32)
                        title_para.font.bold = True
                        title_para.font.color.rgb = RGBColor(0, 51, 102)
                
                content_items = slide_data.get('content', [])
                if content_items and len(slide.shapes.placeholders) > 1:
                    text_frame = slide.shapes.placeholders[1].text_frame
                    text_frame.clear()
                    text_frame.margin_left = Inches(0.2)
                    text_frame.margin_right = Inches(0.2)
                    
                    for i, bullet_point in enumerate(content_items):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = str(bullet_point)
                        p.level = 0
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(51, 51, 51)
                        p.space_after = Pt(16)
                        p.space_before = Pt(8)
        
        return prs
        
    except Exception as e:
        print(f"Error creating presentation: {e}")
        return None

def _create_content_slide(prs, slide_data):
    """Helper function to create a content slide with enhanced formatting"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = slide_data.get('slide_title', 'Content')
    
    title_shape = slide.shapes.title
    if title_shape.has_text_frame:
        title_frame = title_shape.text_frame
        if title_frame.paragraphs:
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    if slide_data.get('content') and len(slide.shapes.placeholders) > 1:
        text_frame = slide.shapes.placeholders[1].text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        
        for i, bullet_point in enumerate(slide_data['content']):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = str(bullet_point)
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.space_after = Pt(16)
            p.space_before = Pt(8)

def create_zip_with_ppts(presentations_data, topic):
    """Create a zip file containing multiple PowerPoint presentations"""
    zip_buffer = io.BytesIO()
    
    try:
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            sessions = presentations_data.get('sessions', [])
            
            for i, session_data in enumerate(sessions, 1):
                try:
                    prs = create_ppt_session(session_data, topic)
                    
                    if prs is None:
                        print(f"Failed to create presentation for session {i}")
                        continue
                    
                    ppt_buffer = io.BytesIO()
                    prs.save(ppt_buffer)
                    ppt_buffer.seek(0)
                    
                    session_title = session_data.get('session_title', f'Session {i}')
                    safe_topic = sanitize_filename(topic)
                    safe_session_title = sanitize_filename(session_title)
                    
                    filename = f"{safe_topic}_Session_{i:02d}_{safe_session_title}.pptx"
                    
                    zip_file.writestr(filename, ppt_buffer.getvalue())
                    ppt_buffer.close()
                    
                    print(f"Created presentation: {filename}")
                    
                except Exception as e:
                    print(f"Error creating presentation for session {i}: {e}")
                    continue
        
        zip_buffer.seek(0)
        return zip_buffer
        
    except Exception as e:
        print(f"Error creating zip file: {e}")
        return None

@app.route("/generate", methods=["POST"])
def generate():
    """Main endpoint to generate PowerPoint presentations"""
    try:
        data = request.json
        topic = data.get("topic", "Course")
        sessions = int(data.get("sessions", 1))
        details = data.get("details", "")
        syllabus = data.get("syllabus", "")
        
        print(f"Generating for topic: {topic}, sessions: {sessions}, details: {details}, syllabus: {syllabus}")

        outline_response = gen_structured_outline(topic, sessions, details, syllabus)
        print("Generated structured outline response:\n", outline_response)
        
        try:
            if "```json" in outline_response:
                json_start = outline_response.find("```json") + 7
                json_end = outline_response.find("```", json_start)
                json_content = outline_response[json_start:json_end].strip()
            elif "```" in outline_response:
                json_start = outline_response.find("```") + 3
                json_end = outline_response.find("```", json_start)
                json_content = outline_response[json_start:json_end].strip()
            else:
                json_content = outline_response.strip()
            
            presentations_data = json.loads(json_content)
            print("Successfully parsed JSON structure")
            
        except json.JSONDecodeError as e:
            print(f"JSON parsing error: {e}")
            print("Falling back to text parsing...")
            presentations_data = parse_text_outline(outline_response, topic, sessions, syllabus)
        
        print(f"Final presentations data structure: {presentations_data}")
        print(f"Number of sessions created: {len(presentations_data.get('sessions', []))}")
        
        if not presentations_data.get('sessions'):
            return jsonify({"error": "No sessions were generated"}), 500
        
        zip_buffer = create_zip_with_ppts(presentations_data, topic)
        
        if zip_buffer is None:
            return jsonify({"error": "Failed to create presentations"}), 500
        
        safe_topic = sanitize_filename(topic)
        zip_filename = f"{safe_topic}_Course_Materials.zip"
        
        return send_file(
            zip_buffer,
            as_attachment=True,
            download_name=zip_filename,
            mimetype="application/zip"
        )
        
    except Exception as e:
        print("❌ Error:", e)
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route("/", methods=["GET"])
def health():
    """Health check endpoint"""
    return jsonify({"status": "ok", "message": "PowerPoint Generator API is running"}), 200

if __name__ == "__main__":
    app.run(host='0.0.0.0', port=5001, debug=True)